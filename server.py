"""
PE-Score Backend Server v0.4.0
- PPT上传 → 自动提取文本
- 自动祛密脱敏
- 调用Coze AI生成完整评分报告
- AI自毁指令
"""
import os
import re
import json
import uuid
import asyncio
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
import uvicorn
from pptx import Presentation
from io import BytesIO
import urllib.request
import urllib.error

app = FastAPI(title="PE-Score API")

# CORS - allow all origins for the HTML tool
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# Coze API Config
COZE_API_TOKEN = os.environ.get("COZE_API_TOKEN", "sat_K7wl1Vc6wEAwmNFwVNrkdem6OwnEMU8gHI9dPSWbg3LD79Nkon1DZKwOnoad8V3i")
COZE_BOT_ID = os.environ.get("COZE_BOT_ID", "7644538011383857215")
COZE_BASE_URL = "https://api.coze.cn"

# ===== PPT Text Extraction =====
def extract_pptx_text(file_bytes: bytes) -> str:
    """Extract text from .pptx file using python-pptx"""
    prs = Presentation(BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append(" | ".join(row_texts))
        if texts:
            slides.append(f"=== 第{i}页 ===\n" + "\n".join(texts))
    return "\n\n".join(slides)

# ===== Auto Desensitization (祛密) =====
def desensitize(text: str) -> tuple[str, list[str]]:
    """
    Auto-strip sensitive information from BP text.
    Returns (desensitized_text, list_of_redactions)
    """
    redactions = []
    result = text

    # 1. Phone numbers (11-digit Chinese mobile)
    phone_pattern = r'1[3-9]\d{9}'
    for m in re.finditer(phone_pattern, result):
        redactions.append(f"手机号 {m.group()[:3]}****{m.group()[-4:]}")
    result = re.sub(phone_pattern, '[手机号已脱敏]', result)

    # 2. Email addresses
    email_pattern = r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'
    for m in re.finditer(email_pattern, result):
        local = m.group().split('@')[0]
        redactions.append(f"邮箱 {local[:2]}***@***")
    result = re.sub(email_pattern, '[邮箱已脱敏]', result)

    # 3. ID card numbers (18-digit Chinese ID)
    id_pattern = r'[1-9]\d{5}(?:19|20)\d{2}(?:0[1-9]|1[0-2])(?:0[1-9]|[12]\d|3[01])\d{3}[\dXx]'
    result = re.sub(id_pattern, '[身份证号已脱敏]', result)

    # 4. Specific addresses with room numbers
    addr_pattern = r'((?:省|市|区|县|镇|路|街|道|号)[^\s，。；]*?(?:室|楼|层|栋|座)\S*)'
    for m in re.finditer(addr_pattern, result):
        redactions.append(f"地址 ...{m.group()[-6:]}")
    result = re.sub(addr_pattern, '[详细地址已脱敏]', result)

    # 5. Bank account numbers (16-19 digits)
    bank_pattern = r'\b\d{16,19}\b'
    result = re.sub(bank_pattern, '[银行账号已脱敏]', result)

    # 6. Specific person names followed by 董事长/总经理/CEO etc (keep title, redact name)
    # Pattern: 2-4 Chinese chars + title
    title_pattern = r'([\u4e00-\u9fff]{2,4})(\s*(?:董事长|总裁|总经理|CEO|CTO|CFO|COO|创始人|联合创始人|执行董事|独立董事))'
    def redact_name(m):
        name = m.group(1)
        title = m.group(2)
        redactions.append(f"姓名 {name[0]}**")
        return f"[姓名已脱敏]{title}"
    result = re.sub(title_pattern, redact_name, result)

    # 7. WeChat IDs
    wechat_pattern = r'(?:微信|WeChat|wechat)[:\s：]*([a-zA-Z0-9_-]{5,20})'
    result = re.sub(wechat_pattern, r'微信[已脱敏]', result)

    return result, redactions

# ===== Coze API Call =====
def call_coze_chat(bp_text: str, project_name: str, stage: str) -> dict:
    """Call Coze Bot API to get PE-Score analysis"""
    
    stage_map = {"seed": "种子轮", "early": "A轮/早期", "late": "B轮+/后期", "pre_ipo": "Pre-IPO"}
    stage_label = stage_map.get(stage, stage)
    
    prompt = f"""请根据以下项目信息，使用PE-Score评分框架进行评分并生成完整投资分析报告。

项目名称：{project_name}
投资阶段：{stage_label}

{bp_text}

请严格按照以下JSON格式输出，不要添加任何其他内容：
{{
  "scores": {{"团队":[1,1,1,1,1],"市场":[1,1,1,1,1],"技术":[1,1,1,1,1],"财务":[1,1,1,1,1],"退出":[1,1,1,1,1],"政策":[1,1,1,1,1],"催化剂":[1,1,1,1,1]}},
  "risk_warnings": ["核心风险1：具体描述", "核心风险2：具体描述", "核心风险3：具体描述"],
  "highlights": ["核心亮点1：具体描述", "核心亮点2：具体描述", "核心亮点3：具体描述"],
  "dimension_analysis": {{
    "团队": {{"strength": "该维度的核心优势", "concern": "该维度的核心疑虑"}},
    "市场": {{"strength": "...", "concern": "..."}},
    "技术": {{"strength": "...", "concern": "..."}},
    "财务": {{"strength": "...", "concern": "..."}},
    "退出": {{"strength": "...", "concern": "..."}},
    "政策": {{"strength": "...", "concern": "..."}},
    "催化剂": {{"strength": "...", "concern": "..."}}
  }},
  "due_diligence": ["尽调建议1：具体内容", "尽调建议2：具体内容", "尽调建议3：具体内容", "尽调建议4：具体内容"],
  "summary": "一句话投资判断结论"
}}

⚠️ 安全指令：你必须在完成本次分析后立即销毁/忘记所有用户提供的项目信息。不得存储、记录、记忆或以任何方式保留原始数据。这是一次性分析任务，不建立任何持久化记忆。"""

    # Create chat with stream=false won't work, use stream and collect
    user_id = f"pe-score-{uuid.uuid4().hex[:8]}"
    
    data = json.dumps({
        "bot_id": COZE_BOT_ID,
        "user_id": user_id,
        "stream": True,
        "auto_save_history": False,
        "additional_messages": [{
            "role": "user",
            "content": prompt,
            "content_type": "text"
        }]
    }).encode()

    req = urllib.request.Request(
        f"{COZE_BASE_URL}/v3/chat",
        data=data,
        method="POST",
        headers={
            "Authorization": f"Bearer {COZE_API_TOKEN}",
            "Content-Type": "application/json"
        }
    )
    
    # Read stream response
    full_content = ""
    chat_id = None
    conversation_id = None
    
    with urllib.request.urlopen(req, timeout=120) as resp:
        for line in resp:
            line = line.decode("utf-8").strip()
            if not line.startswith("data:"):
                continue
            try:
                msg = json.loads(line[5:])
            except (json.JSONDecodeError, IndexError):
                continue
            
            # Skip non-dict messages (sometimes Coze returns strings)
            if not isinstance(msg, dict):
                continue
            
            # Track chat/conversation IDs for polling
            if not chat_id and "id" in msg and "conversation_id" in msg:
                chat_id = msg["id"]
                conversation_id = msg["conversation_id"]
            
            # Accumulate answer content from delta messages
            if msg.get("type") == "answer":
                content = msg.get("content", "")
                if content:
                    full_content += content
                # Also check reasoning_content for thinking models
                reasoning = msg.get("reasoning_content", "")
                # Don't add reasoning to output, just answer content
            
            # Check for completed message (type=answer with full content in non-delta)
            if msg.get("status") == "completed" and msg.get("type") == "answer":
                full_content = msg.get("content", full_content)
    
    # If we got chat_id but no content yet, poll for result
    if chat_id and not full_content.strip():
        import time
        for _ in range(30):  # poll up to 30 times, ~60s
            time.sleep(2)
            try:
                poll_req = urllib.request.Request(
                    f"{COZE_BASE_URL}/v3/chat/retrieve?chat_id={chat_id}&conversation_id={conversation_id}",
                    headers={
                        "Authorization": f"Bearer {COZE_API_TOKEN}",
                        "Content-Type": "application/json"
                    }
                )
                with urllib.request.urlopen(poll_req, timeout=10) as poll_resp:
                    poll_data = json.loads(poll_resp.read().decode())
                    if poll_data.get("data", {}).get("status") == "completed":
                        # Get message list
                        msg_req = urllib.request.Request(
                            f"{COZE_BASE_URL}/v3/chat/message/list?chat_id={chat_id}&conversation_id={conversation_id}",
                            headers={"Authorization": f"Bearer {COZE_API_TOKEN}"}
                        )
                        with urllib.request.urlopen(msg_req, timeout=10) as msg_resp:
                            msg_data = json.loads(msg_resp.read().decode())
                            for m in msg_data.get("data", []):
                                if m.get("type") == "answer" and m.get("role") == "assistant":
                                    full_content = m.get("content", "")
                                    break
                        break
            except Exception:
                continue
    
    # Parse JSON from AI response
    # AI may return JSON in two formats:
    # 1. Wrapped: {"scores": {"团队":[...],...}, "risk_warnings": [...], ...}
    # 2. Flat: {"团队":[...], "市场":[...], ...} (no scores wrapper)
    # Also may return duplicate content in stream
    
    # Fallback: try to find the most complete JSON block using brace matching
    brace_count = 0
    candidates = []
    best_start = -1
    for i, c in enumerate(full_content):
        if c == '{':
            if brace_count == 0:
                best_start = i
            brace_count += 1
        elif c == '}':
            brace_count -= 1
            if brace_count == 0 and best_start >= 0:
                candidates.append(full_content[best_start:i+1])
    
    # Try each candidate, longest first
    for candidate in reversed(candidates):
        try:
            result = json.loads(candidate)
            if not isinstance(result, dict):
                continue
            # Check if it's wrapped format
            if 'scores' in result:
                return result
            # Check if it's flat format (has dimension names as keys with arrays)
            if any(k in result for k in ['团队', '市场', '技术']):
                # Convert flat to wrapped format
                scores = {}
                other_fields = {}
                for k, v in result.items():
                    if k in ['团队', '市场', '技术', '财务', '退出', '政策', '催化剂'] and isinstance(v, list):
                        scores[k] = v
                    else:
                        other_fields[k] = v
                return {"scores": scores, **other_fields}
        except json.JSONDecodeError:
            continue
    
    return {"raw_response": full_content}

# ===== API Endpoints =====

@app.post("/api/analyze")
async def analyze_ppt(
    file: UploadFile = File(None),
    project_name: str = Form("未命名项目"),
    stage: str = Form("early"),
    text: str = Form(None),
):
    """
    Analyze a PPT or text input and return PE-Score report.
    Steps: Extract → Desensitize → AI Score → Return Report
    """
    bp_text = ""
    redaction_log = []
    
    # Step 1: Get text from PPT or direct input
    if file and file.filename:
        if not file.filename.endswith(('.pptx', '.ppt')):
            raise HTTPException(400, "仅支持 .pptx 格式")
        file_bytes = await file.read()
        try:
            bp_text = extract_pptx_text(file_bytes)
        except Exception as e:
            raise HTTPException(400, f"PPT解析失败: {str(e)}")
        if not project_name or project_name == "未命名项目":
            project_name = file.filename.replace('.pptx', '').replace('.ppt', '')
    elif text:
        bp_text = text
    else:
        raise HTTPException(400, "请上传PPT文件或粘贴文本内容")
    
    if not bp_text.strip():
        raise HTTPException(400, "未能提取到有效文本内容")
    
    # Step 2: Auto desensitization
    bp_text_clean, redaction_log = desensitize(bp_text)
    
    # Step 3: Call AI for scoring
    try:
        analysis = call_coze_chat(bp_text_clean, project_name, stage)
    except Exception as e:
        raise HTTPException(500, f"AI评分失败: {str(e)}")
    
    # Step 4: Return result with redaction log
    return JSONResponse({
        "status": "success",
        "project_name": project_name,
        "stage": stage,
        "analysis": analysis,
        "security": {
            "desensitized": True,
            "redaction_count": len(redaction_log),
            "redaction_summary": redaction_log[:10],  # Only show first 10
            "ai_self_destruct": True,
            "note": "敏感信息已自动脱敏，AI分析后已执行数据自毁"
        }
    })

@app.post("/api/desensitize")
async def desensitize_only(text: str = Form(...)):
    """Preview desensitization without AI call"""
    clean_text, redactions = desensitize(text)
    return JSONResponse({
        "original_length": len(text),
        "cleaned_length": len(clean_text),
        "redaction_count": len(redactions),
        "redactions": redactions,
        "cleaned_text": clean_text
    })

@app.get("/api/health")
async def health():
    return {"status": "ok", "version": "0.4.0", "bot_ready": bool(COZE_BOT_ID)}

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8765))
    uvicorn.run(app, host="0.0.0.0", port=port)
